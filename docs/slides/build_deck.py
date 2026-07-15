#!/usr/bin/env python3
"""Build the OX.ar Embryo Module "Idiot's Guide" implementation deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- OX.ar brand palette ----
INK   = RGBColor(0x22, 0x1f, 0x20)
AQUA  = RGBColor(0x45, 0xa5, 0x9d)
DEEP  = RGBColor(0x2f, 0x7d, 0x76)
SOFT  = RGBColor(0xea, 0xf5, 0xf3)
PAPER = RGBColor(0xf5, 0xf8, 0xf7)
WHITE = RGBColor(0xff, 0xff, 0xff)
MUTED = RGBColor(0x6c, 0x7c, 0x7a)
GREEN = RGBColor(0x2f, 0x9e, 0x6f)
AMBER = RGBColor(0xc9, 0x82, 0x1c)
RED   = RGBColor(0xcf, 0x4a, 0x4a)
PURPLE= RGBColor(0x7a, 0x5c, 0xc0)
LINE  = RGBColor(0xe4, 0xec, 0xea)

HEAD = "Montserrat"
BODY = "Lato"
MONO = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.06):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, font, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = font
            if rest and rest[0]:
                r.font.italic = True
    return tb


def header(slide, eyebrow, title, n):
    box(slide, 0, 0, EMU_W, Inches(1.15), fill=WHITE)
    box(slide, 0, Inches(1.15), EMU_W, Pt(3), fill=AQUA)
    text(slide, Inches(0.6), Inches(0.16), Inches(11), Inches(0.4),
         [[(eyebrow, 12, DEEP, True, HEAD)]])
    text(slide, Inches(0.6), Inches(0.44), Inches(11.5), Inches(0.7),
         [[(title, 27, INK, True, HEAD)]])
    # slide number chip
    c = box(slide, Inches(12.35), Inches(0.32), Inches(0.62), Inches(0.5),
            fill=SOFT, rounded=True)
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = DEEP; r.font.name = HEAD


def bullet(slide, x, y, w, items, gap=0.5, size=15, num=True):
    """items: list of (bold_lead, rest) ; renders numbered/dotted rows."""
    for i, (lead, rest) in enumerate(items):
        yy = y + Inches(gap * i)
        if num:
            d = box(slide, x, yy + Inches(0.02), Inches(0.34), Inches(0.34),
                    fill=AQUA, rounded=True)
            tf = d.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(i + 1); r.font.size = Pt(13)
            r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HEAD
            tx = x + Inches(0.5)
        else:
            box(slide, x, yy + Inches(0.12), Inches(0.14), Inches(0.14), fill=DEEP)
            tx = x + Inches(0.34)
        runs = [[(lead, size, INK, True, BODY)] + ([(rest, size, MUTED, False, BODY)] if rest else [])]
        text(slide, tx, yy - Inches(0.02), w - (tx - x), Inches(0.5), runs, line_spacing=1.05)


def chip(slide, x, y, w, h, label, fill, txtcol=WHITE, size=12, font=HEAD, bold=True):
    s = box(slide, x, y, w, h, fill=fill, rounded=True)
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = txtcol; r.font.name = font
    return s


# ============================================================== SLIDE 1 — TITLE
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 0, 0, EMU_W, Inches(2.5), fill=INK)
box(s, 0, Inches(4.55), EMU_W, Pt(4), fill=AQUA)
# logo
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(8), Inches(1.2))
tf = tb.text_frame
p = tf.paragraphs[0]
for t, c in [("OX", WHITE), (".", AQUA), ("ar", WHITE)]:
    r = p.add_run(); r.text = t; r.font.size = Pt(54); r.font.bold = True
    r.font.name = HEAD; r.font.color.rgb = c
text(s, Inches(0.95), Inches(2.7), Inches(11), Inches(1.4),
     [[("The Idiot's Guide to Building the", 20, RGBColor(0xcf,0xe6,0xe2), False, BODY)],
      [("Embryo Module", 40, WHITE, True, HEAD)]], space_after=2)
text(s, Inches(0.95), Inches(4.8), Inches(11.6), Inches(2),
     [[("A dead-simple, step-by-step plan for OX.assisted reproduction (OX.ar)", 16, RGBColor(0xcf,0xe6,0xe2), False, BODY)],
      [("Module M-18  ·  Embryology Cycle Management", 14, AQUA, True, HEAD)],
      [("Repos: OxArFrontendReact (React)  +  OxArBackendReact (.NET Web API + Dataverse)", 12.5, MUTED, False, MONO)]],
     space_after=8)
chip(s, Inches(0.95), Inches(6.55), Inches(3.7), Inches(0.5),
     "100% SEPARATE FROM OX.gp", RED, size=12)

# ====================================================== SLIDE 2 — WHAT / WHY
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "START HERE", "What are we building, in plain English?", 2)
text(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.2),
     [[("It's the ", 16, INK, False, BODY), ("lab workspace embryologists and clinicians use", 16, DEEP, True, BODY),
       (" to run an IVF cycle: plan the day's events, record each procedure, grade embryos day by day, "
        "manage frozen storage, and keep the safety/witnessing record the law requires.", 16, INK, False, BODY)]])
cards = [
    ("PLAN", "Worklist, cycle overview, consent & 3-point check", AQUA),
    ("DO", "Egg collection, semen prep, insemination, fertilisation", DEEP),
    ("GROW", "Day 0-7 embryo grading, then transfer", GREEN),
    ("STORE", "Freeze / thaw, biopsy, cryo tanks & inventory", PURPLE),
]
cw = Inches(2.95); gap = Inches(0.15); x0 = Inches(0.6); y0 = Inches(3.0)
for i, (t, d, col) in enumerate(cards):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, Inches(2.5), fill=WHITE, line=LINE, rounded=True)
    box(s, x, y0, cw, Inches(0.12), fill=col)
    chip(s, x + Inches(0.3), y0 + Inches(0.35), Inches(1.3), Inches(0.5), t, col, size=15)
    text(s, x + Inches(0.3), y0 + Inches(1.05), cw - Inches(0.6), Inches(1.4),
         [[(d, 13.5, INK, False, BODY)]])
text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(1),
     [[("The spec = 41 requirements (M18-000 … M18-040 + 1 new) across 23 data sheets. "
        "This deck tells you how to build it against the REAL OX.ar code.", 13.5, MUTED, False, BODY)]])

# ====================================================== SLIDE 3 — GOLDEN RULES
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "READ TWICE", "3 golden rules before you touch anything", 3)
rules = [
    ("Never mix with OX.gp.", " No shared code, entities, endpoints, assets or contracts. "
     "Only the two OX.ar repos + the OX.ar Dataverse org. If a task seems to need OX.gp — STOP."),
    ("Don't rebuild what exists.", " Dataverse already has bcrm_egg (the specimen) and bcrm_eggdetail "
     "(its day-by-day assessments). Reuse them. Add only the 15 missing workflow tables."),
    ("Witnessing is not optional.", " Every critical step needs an independent second person. "
     "The app must BLOCK a procedure until a 3-point check passes. This is the law (HFEA)."),
]
bullet(s, Inches(0.7), Inches(1.7), Inches(12), rules, gap=1.35, size=17)
chip(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.62),
     "One specimen identity per egg/embryo, cradle to grave = safe traceability", SOFT, DEEP, size=14)

# ====================================================== SLIDE 4 — BIG PICTURE
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "THE MAP", "How the pieces fit together", 4)
layers = [
    ("REACT FRONTEND", "OxArFrontendReact  ·  a new /lab area (screens for the lab team)", AQUA),
    (".NET WEB API", "OxArBackendReact  ·  Controller → Service → CRMCoreRepository", DEEP),
    ("DATAVERSE", "contact · bcrm_treatment_cycle · bcrm_egg · bcrm_eggdetail · +15 new tables", INK),
]
y = Inches(1.7)
for i, (t, d, col) in enumerate(layers):
    box(s, Inches(1.6), y, Inches(10.1), Inches(1.15), fill=WHITE, line=LINE, rounded=True)
    box(s, Inches(1.6), y, Inches(0.16), Inches(1.15), fill=col)
    chip(s, Inches(1.95), y + Inches(0.32), Inches(2.9), Inches(0.5), t, col, size=13)
    text(s, Inches(5.1), y + Inches(0.18), Inches(6.4), Inches(0.9),
         [[(d, 13.5, INK, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.35), y + Inches(1.16),
                               Inches(0.55), Inches(0.32))
        a.fill.solid(); a.fill.fore_color.rgb = AQUA; a.line.fill.background(); a.shadow.inherit = False
    y += Inches(1.63)
text(s, Inches(1.6), Inches(6.7), Inches(10), Inches(0.6),
     [[("Auth already exists: MSAL Azure AD B2C + the hmacauth token. Reuse it — build no new login.", 13, MUTED, False, BODY)]])

# ====================================================== SLIDE 5 — BUILD ORDER
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "THE RHYTHM", "The order you build EVERY feature", 5)
steps = ["Dataverse\ntable", "API\nendpoint", "React\nservice", "React\nscreen", "Safety\ngate", "Verify\nend-to-end"]
cols = [INK, DEEP, AQUA, AQUA, RED, GREEN]
x = Inches(0.55); y = Inches(2.4); w = Inches(1.85); gap = Inches(0.18)
for i, (st, col) in enumerate(zip(steps, cols)):
    xx = x + i * (w + gap)
    box(s, xx, y, w, Inches(1.7), fill=WHITE, line=col, line_w=2, rounded=True)
    chip(s, xx + Inches(0.65), y + Inches(0.2), Inches(0.55), Inches(0.55), str(i+1), col, size=16)
    text(s, xx + Inches(0.1), y + Inches(0.85), w - Inches(0.2), Inches(0.8),
         [[(l, 14, INK, True, HEAD)] for l in st.split("\n")], align=PP_ALIGN.CENTER, space_after=0)
    if i < len(steps) - 1:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, xx + w + Inches(0.01),
                               y + Inches(0.72), Inches(0.16), Inches(0.26))
        a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background(); a.shadow.inherit = False
text(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(1),
     [[("Rule: ", 16, RED, True, BODY),
       ("never start a screen before its table and endpoint exist. No mock data in real screens.",
        16, INK, False, BODY)]])
text(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.2),
     [[("Do the whole rhythm for one small slice, demo it, then move to the next. "
        "That's how the 5 phases (next slide) stay shippable.", 13.5, MUTED, False, BODY)]])

# ====================================================== SLIDE 6 — PHASES
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "THE PLAN", "5 phases — each one is demoable on its own", 6)
phases = [
    ("0", "Foundations", "Confirm schema · create tables & option sets · build the /lab shell + witness chip", INK),
    ("1", "Planning", "Worklist · cycle overview · consent + 3-point-check GATE working end-to-end", AQUA),
    ("2", "Procedures", "Egg collection (MII/MI/GV) · semen prep · insemination · fertilisation (PN board)", DEEP),
    ("3", "Lab Culture", "The Day 0-7 development grid (Gardner grades) · transfer · live KPIs", GREEN),
    ("4", "Cryo & Governance", "Tanks & inventory · freeze/thaw/biopsy · witnessing log · research banner", PURPLE),
    ("5", "Reports & Manager", "Outcome reports & letters · lab worksheet · manager correct/delete (audited)", AMBER),
]
y = Inches(1.55)
for num, name, desc, col in phases:
    box(s, Inches(0.6), y, Inches(12.1), Inches(0.82), fill=WHITE, line=LINE, rounded=True)
    chip(s, Inches(0.78), y + Inches(0.16), Inches(0.5), Inches(0.5), num, col, size=15)
    text(s, Inches(1.5), y + Inches(0.1), Inches(2.5), Inches(0.62),
         [[(name, 15, INK, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.0), y + Inches(0.1), Inches(8.5), Inches(0.62),
         [[(desc, 12.5, MUTED, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.9)

# ====================================================== SLIDE 7 — BACKEND RECIPE
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "BACKEND RECIPE", "Adding one feature = 4 files (copy an existing one)", 7)
files = [
    ("1. DTO", "TaskService.DataTransferObjects/<Area>/XxxDTO.cs", "plain properties (copy StoredSampleDTO)"),
    ("2. Service", "TaskService.Services/<Area>/IXxxService.cs + XxxService.cs", "reads = FetchXML · writes = Entity"),
    ("3. Controller", "TaskService/Controllers/XxxController.cs", "[Authorize] + [RoutePrefix] · copy an existing one"),
    ("4. (Repository)", "no change", "reuse CRMCoreRepository.Create/Update/Delete/Get"),
]
y = Inches(1.55)
for lead, path, note in files:
    box(s, Inches(0.6), y, Inches(12.1), Inches(0.92), fill=WHITE, line=LINE, rounded=True)
    text(s, Inches(0.85), y + Inches(0.13), Inches(2.4), Inches(0.7),
         [[(lead, 15, DEEP, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.15), y + Inches(0.11), Inches(9.3), Inches(0.75),
         [[(path, 12, INK, True, MONO)],
          [(note, 12, MUTED, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    y += Inches(1.0)
text(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.2),
     [[("No dependency-injection wiring to touch. ", 14, INK, True, BODY),
       ("Attribute routing means a new controller just works. Test it in Swagger first.",
        14, MUTED, False, BODY)]])
text(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.6),
     [[("Idioms: OptionSetValue(int) · new EntityReference(\"logical\", guid) · SpecifyKind(date, Utc) · run the ACSHelper guard.",
        11.5, DEEP, False, MONO)]])

# ====================================================== SLIDE 8 — FRONTEND RECIPE
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "FRONTEND RECIPE", "Adding one screen", 8)
fe = [
    ("Put a route", " in src/Layout Component/Layout1.js under a new /lab tree (staff-only)."),
    ("Add API calls", " to one file: src/services/embryoApi.js (axios, same style as ProfileService.js)."),
    ("Build the page", " in src/lab/pages/… — load in useEffect, write via embryoApi, toast the result."),
    ("Reuse components", " WitnessChip · StatusBadge · StatTile · PnBoard · DevelopmentGrid."),
    ("Show the gate", " disable the button until the 3-point check passes & witness ≠ performer."),
]
bullet(s, Inches(0.7), Inches(1.7), Inches(12), fe, gap=0.82, size=15.5)
chip(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7),
     "Keep ALL lab code in src/lab/ + embryoApi.js. Don't restyle the patient portal.", SOFT, DEEP, size=13.5)

# ====================================================== SLIDE 9 — FULL CRUD
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "IT'S NOT READ-ONLY", "Everything can be added, edited, updated & deleted", 9)
text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.75),
     [[("Every record ", 16, INK, True, BODY),
       ("— and every list that drives the forms (templates, tanks, sites, research projects) — is ",
        16, INK, False, BODY),
       ("fully read-write from the UI.", 16, DEEP, True, BODY)]])
crud = [
    ("ADD", "Create new rows. Every subgrid has an + Add button.", AQUA),
    ("EDIT", "Partial update — send only what changed (incl. inline grid cells).", DEEP),
    ("DELETE", "Soft-retire with a reason by default; hard delete = manager only.", RED),
    ("MANAGE LISTS", "Templates, cryo tanks, research projects — CRUD, no developer needed.", PURPLE),
]
x0 = Inches(0.6); y0 = Inches(2.35); cw = Inches(2.95); gap = Inches(0.15)
for i, (t, d, col) in enumerate(crud):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, Inches(2.15), fill=WHITE, line=LINE, rounded=True)
    box(s, x, y0, cw, Inches(0.12), fill=col)
    chip(s, x + Inches(0.3), y0 + Inches(0.35), Inches(2.0), Inches(0.5), t, col, size=13)
    text(s, x + Inches(0.3), y0 + Inches(1.05), cw - Inches(0.6), Inches(1.0),
         [[(d, 12.5, INK, False, BODY)]])
box(s, Inches(0.6), Inches(4.95), Inches(12.1), Inches(1.5), fill=SOFT, rounded=True)
text(s, Inches(0.85), Inches(5.08), Inches(11.6), Inches(0.45),
     [[("Guardrails when editing clinical data:", 14, DEEP, True, HEAD)]])
guards = [
    ("Concurrency", "two people, one cycle → 409 on a clash, reload"),
    ("Audit trail", "Dataverse logs every field edit"),
    ("Role-gated", "view / edit / delete by role (403 on server)"),
    ("Correct & re-enter", "managers fix locked data with a reason"),
]
gx = Inches(0.85); gw = Inches(2.85)
for i, (t, d) in enumerate(guards):
    x = gx + i * (gw + Inches(0.1))
    text(s, x, Inches(5.55), gw, Inches(0.8),
         [[(t, 12.5, INK, True, BODY)], [(d, 11, MUTED, False, BODY)]], space_after=1)

# ====================================================== SLIDE 10 — SAFETY GATES
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "NON-NEGOTIABLE", "The safety gates (enforce on the SERVER too)", 10)
gates = [
    ("Independent witness", "reject if witness = performer", RED),
    ("3-point check", "no procedure starts until it PASSES", RED),
    ("Consent gate", "cycle must be 'ready for procedure'", AMBER),
    ("ICSI = MII only", "only mature (MII) eggs go to ICSI", DEEP),
    ("2PN = normal", "flag 0PN / 1PN / 3PN as abnormal", AMBER),
    ("Carry-over", "un-fated embryos roll to next day", GREEN),
]
x0 = Inches(0.6); y0 = Inches(1.6); cw = Inches(3.9); ch = Inches(1.5); gx = Inches(0.2); gy = Inches(0.25)
for i, (t, d, col) in enumerate(gates):
    r, c = divmod(i, 3)
    x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
    box(s, x, y, cw, ch, fill=WHITE, line=LINE, rounded=True)
    box(s, x, y, cw, Inches(0.1), fill=col)
    text(s, x + Inches(0.25), y + Inches(0.28), cw - Inches(0.5), Inches(0.5),
         [[(t, 16, INK, True, HEAD)]])
    text(s, x + Inches(0.25), y + Inches(0.82), cw - Inches(0.5), Inches(0.6),
         [[(d, 13, MUTED, False, BODY)]])
text(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(1),
     [[("Why server-side? ", 14, RED, True, BODY),
       ("The UI can be bypassed. The API is the last line of defence for patient safety.",
        14, INK, False, BODY)]])

# ====================================================== SLIDE 11 — CLINICAL CHEATSHEET
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "CHEAT SHEET", "The clinical words you must get right", 11)
terms = [
    ("MII / MI / GV", "egg maturity. Only MII (mature) can be ICSI-inseminated."),
    ("PN (pronuclei)", "checked ~16-18h after insemination. 2PN = normal fertilisation."),
    ("PB (polar body)", "recorded alongside PN at the fertilisation check."),
    ("Gardner grade", "blastocyst = expansion 1-6 + ICM (A/B/C) + TE (A/B/C), e.g. 4AA."),
    ("Fragmentation %", "cleavage-stage quality; high = worse."),
    ("Vienna KPIs", "benchmarks, e.g. Day-5 blastocyst rate 44-80% of 2PN zygotes."),
    ("Istanbul 2025", "the current ESHRE/ALPHA grading standard the terminology follows."),
    ("Double witnessing", "independent 2nd person on every critical step (HFEA)."),
]
y = Inches(1.55)
for i, (t, d) in enumerate(terms):
    box(s, Inches(0.6), y, Inches(12.1), Inches(0.62), fill=WHITE if i % 2 == 0 else SOFT, rounded=False)
    text(s, Inches(0.8), y + Inches(0.08), Inches(3.2), Inches(0.5),
         [[(t, 13.5, DEEP, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.1), y + Inches(0.08), Inches(8.4), Inches(0.5),
         [[(d, 12.5, INK, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.64)

# ====================================================== SLIDE 12 — DESIGN
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "MAKE IT LOOK RIGHT", "Branding cheat sheet (match the prototype)", 12)
text(s, Inches(0.6), Inches(1.45), Inches(12), Inches(0.6),
     [[("Fonts:  ", 15, INK, True, BODY), ("Montserrat", 15, DEEP, True, HEAD),
       (" headings · ", 14, MUTED, False, BODY), ("Lato", 15, DEEP, True, BODY),
       (" body · ", 14, MUTED, False, BODY), ("JetBrains Mono", 15, DEEP, True, MONO),
       (" for IDs & lab numbers", 14, MUTED, False, BODY)]])
swatches = [
    ("#45a59d", "aqua", AQUA, WHITE), ("#2f7d76", "buttons", DEEP, WHITE),
    ("#2f9e6f", "normal / pass", GREEN, WHITE), ("#c9821c", "attention", AMBER, WHITE),
    ("#cf4a4a", "abnormal / fail", RED, WHITE), ("#7a5cc0", "research", PURPLE, WHITE),
]
x0 = Inches(0.6); y0 = Inches(2.35); cw = Inches(1.95); gx = Inches(0.05)
for i, (hexv, name, col, tc) in enumerate(swatches):
    x = x0 + i * (cw + gx)
    box(s, x, y0, cw, Inches(1.5), fill=col, rounded=True)
    text(s, x, y0 + Inches(0.45), cw, Inches(0.4), [[(hexv, 13, tc, True, MONO)]], align=PP_ALIGN.CENTER)
    text(s, x, y0 + Inches(0.85), cw, Inches(0.4), [[(name, 11.5, tc, False, BODY)]], align=PP_ALIGN.CENTER)
rules2 = [
    ("Colour is never the only signal", " — always pair it with a label or icon (safety-critical, colour-blind-safe)."),
    ("Every specimen ID & lab value", " renders in monospace so numbers line up."),
    ("All demo data stays synthetic", " — SAMPLE / MODEL / DEMO / PROTO patients only."),
]
bullet(s, Inches(0.7), Inches(4.3), Inches(12), rules2, gap=0.72, size=14.5, num=False)

# ====================================================== SLIDE 13 — DONE + NEXT
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "FINISH LINE", "\"Done\" checklist  +  where everything lives", 13)
done = [
    "Table exists in the OXAR_EmbryoModule solution (option sets + relationships)",
    "Endpoint smoke-tested in Swagger; lookups & option sets round-trip",
    "React screen wired to the LIVE endpoint (no mock data)",
    "Witness / consent / eligibility rules enforced server-side AND shown in UI",
    "Brand matches the prototype; IDs & values in mono",
    "Requirement ticked; MIVF / SME notified for clinical sign-off",
]
box(s, Inches(0.6), Inches(1.5), Inches(7.2), Inches(4.9), fill=WHITE, line=LINE, rounded=True)
text(s, Inches(0.85), Inches(1.65), Inches(6.7), Inches(0.5), [[("Definition of Done (every slice)", 15, DEEP, True, HEAD)]])
y = Inches(2.25)
for d in done:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.04), Inches(0.28), Inches(0.28))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN; c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓"; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HEAD
    text(s, Inches(1.35), y - Inches(0.02), Inches(6.3), Inches(0.7), [[(d, 12.5, INK, False, BODY)]])
    y += Inches(0.7)
# right column — docs
box(s, Inches(8.0), Inches(1.5), Inches(4.7), Inches(4.9), fill=INK, rounded=True)
text(s, Inches(8.3), Inches(1.7), Inches(4.2), Inches(0.5), [[("What's in the repo", 15, AQUA, True, HEAD)]])
items = [
    ("docs/", "README + 01-06: data model, API, frontend, CRUD, code-drop, treatment-type templates"),
    ("backend/", "working C#: EmbryoEvent, Opu, CryoTank"),
    ("docs/templates/", "treatment-type-display-templates.json"),
    ("docs/slides/", "this deck + generator"),
]
y = Inches(2.35)
for f, d in items:
    text(s, Inches(8.3), y, Inches(4.2), Inches(0.75),
         [[(f, 12.5, WHITE, True, MONO)], [(d, 10.5, RGBColor(0xcf,0xe6,0xe2), False, BODY)]], space_after=1)
    y += Inches(0.82)
text(s, Inches(8.3), Inches(5.9), Inches(4.2), Inches(0.6),
     [[("github.com/Jkosobucki/OXARv1", 10.5, AQUA, False, MONO)]], space_after=2)

# ====================================================== SLIDE 14 — TEMPLATES CONCEPT
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "SHOW ONLY WHAT MATTERS", "One template per treatment type", 14)
text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.0),
     [[("Today every cycle shows every tab and field. Instead, the ", 16, INK, False, BODY),
       ("treatment type", 16, DEEP, True, BODY),
       (" (with the patient) drives a template that ", 16, INK, False, BODY),
       ("hides what's irrelevant and surfaces the right steps", 16, DEEP, True, BODY),
       (" — faster, safer data entry for the embryologist.", 16, INK, False, BODY)]])
how = [
    ("Pick patient + treatment type", " on the Treatment Cycle."),
    ("Template loads", " — hides the tabs/fields that don't apply, shows the ones that do."),
    ("Lab steps seed automatically", " — the worklist lists the right procedures for that type."),
    ("Managers can tune templates", " in the admin screen (reference data, not hard-coded)."),
]
bullet(s, Inches(0.7), Inches(2.7), Inches(12), how, gap=0.62, size=15)
chip(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(0.7),
     "Reuses the Event Template design (bcrm_eventtemplate keyed by treatment type)", SOFT, DEEP, size=13.5)
text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
     [[("8 treatment types: IVF · Egg Only · Embryo Transfer · VOT · IUI · Ovulation Induction · Tracking · Thaw & Refreeze",
        12.5, MUTED, False, BODY)]])

# ====================================================== SLIDE 15 — MATRIX
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "AT A GLANCE", "Key lab steps by treatment type", 15)
cols = ["IVF", "Egg", "FET", "VOT", "IUI", "OI", "Trk", "T+R"]
rows = [
    ("Consent & 3-Point Check", ["Y","Y","Y","Y","Y","Y","-","Y"]),
    ("Ovarian Stimulation (Drugs)", ["Y","Y","-","-","~","Y","-","-"]),
    ("Folliculogram monitoring", ["Y","Y","~","-","Y","Y","Y","-"]),
    ("OPU (Egg Collection)", ["Y","Y","-","-","-","-","-","-"]),
    ("Oocyte Thaw", ["-","-","-","Y","-","-","-","Y"]),
    ("Semen Prep", ["Y","-","-","Y","Y","-","-","-"]),
    ("Insemination (IVF/ICSI)", ["Y","-","-","Y","-","-","-","-"]),
    ("IUI insemination", ["-","-","-","-","Y","-","-","-"]),
    ("Fertilisation + Culture", ["Y","-","-","Y","-","-","-","-"]),
    ("Embryo Thaw", ["~","-","Y","-","-","-","-","~"]),
    ("Embryo Transfer", ["Y","-","Y","Y","-","-","-","-"]),
    ("Freeze / Store", ["Y","Y","-","~","-","-","-","Y"]),
    ("Outcome / Preg test", ["Y","-","Y","Y","Y","Y","~","-"]),
]
lx = Inches(0.6); lw = Inches(3.3); cw = Inches(1.15); top = Inches(1.55); rh = Inches(0.37)
for j, c in enumerate(cols):
    chip(s, lx + lw + j * cw, top, cw - Inches(0.06), Inches(0.34), c, DEEP, size=10)
for i, (label, vals) in enumerate(rows):
    yy = top + Inches(0.46) + i * rh
    if i % 2 == 0:
        box(s, lx, yy, lw + 8 * cw - Inches(0.06), rh - Inches(0.04), fill=SOFT)
    text(s, lx + Inches(0.06), yy, lw, rh, [[(label, 10.5, INK, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE)
    for j, v in enumerate(vals):
        glyph = {"Y": "✓", "~": "––", "-": "·"}[v]
        gcol = {"Y": GREEN, "~": AMBER, "-": MUTED}[v]
        text(s, lx + lw + j * cw, yy, cw, rh, [[(glyph, 12, gcol, True, HEAD)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, lx, top + Inches(0.46) + len(rows) * rh + Inches(0.05), Inches(12), Inches(0.4),
     [[("✓ applies      –– sometimes      · not applicable", 11, MUTED, False, BODY)]])

# ====================================================== SLIDES 16-23 — PER TYPE
def chip_flow(sl, x0, y0, maxx, steps, fill):
    x = x0; y = y0; lh = Inches(0.36)
    for st in steps:
        w = Inches(max(0.95, 0.115 * len(st) + 0.34))
        if x + w > maxx:
            x = x0; y += lh + Inches(0.14)
        chip(sl, x, y, w, lh, st, fill, WHITE, size=10.5)
        x += w + Inches(0.12)
    return y + lh

TYPES = [
    {"name": "IVF", "color": DEEP,
     "purpose": "Full stimulated cycle — collect eggs, fertilise, transfer / freeze",
     "steps": ["Consent & 3PC", "Stimulation", "Folliculogram", "Trigger", "OPU", "Semen Prep",
               "Insemination", "Fertilisation", "Culture D2-7", "Transfer", "Freeze", "Outcome"],
     "show": ["Overview", "Folliculogram", "Drugs", "Oocyte Retrieval", "Gametes", "Lab Results",
              "Embryo Transfer", "Stored Sample", "Outcome Data", "Checklists"],
     "hide": ["IUI/DI"],
     "fields": "Show partner, diagnosis, sperm source, # embryos to replace"},
    {"name": "Egg Only", "color": GREEN,
     "purpose": "Egg freezing / preservation — collect & vitrify (no fertilisation)",
     "steps": ["Consent & 3PC", "Stimulation", "Folliculogram", "Trigger", "OPU",
               "Denuding / MII", "Vitrify Oocytes", "Storage"],
     "show": ["Overview", "Folliculogram", "Drugs", "Oocyte Retrieval", "Gametes", "Stored Sample", "Checklists"],
     "hide": ["Embryo Transfer", "IUI/DI", "Lab Results", "Recommendation", "Outcome Data"],
     "fields": "Hide partner, sperm source, # embryos to replace"},
    {"name": "Embryo Transfer (FET)", "color": PURPLE,
     "purpose": "Frozen embryo transfer — thaw a stored embryo & transfer",
     "steps": ["Consent & 3PC", "Endo Prep / Lining", "Select Embryo", "Embryo Thaw", "3PC",
               "Transfer", "Outcome"],
     "show": ["Overview", "Folliculogram", "Drugs", "Gametes", "Stored Sample", "Embryo Transfer",
              "Outcome Data", "Checklists"],
     "hide": ["Oocyte Retrieval", "IUI/DI", "Lab Results"],
     "fields": "Show # embryos to replace; hide sperm source & egg-collection fields"},
    {"name": "VOT  (Vitrified Oocyte Thaw — confirm)", "color": AMBER,
     "purpose": "Thaw frozen eggs — ICSI — culture — transfer",
     "steps": ["Consent & 3PC", "Select Frozen Eggs", "Oocyte Thaw", "Survival", "Semen Prep",
               "ICSI", "Fertilisation", "Culture", "Transfer / Freeze", "Outcome"],
     "show": ["Overview", "Gametes", "Stored Sample", "Lab Results", "Embryo Transfer", "Drugs",
              "Outcome Data", "Checklists"],
     "hide": ["Oocyte Retrieval", "Folliculogram", "IUI/DI"],
     "fields": "Confirm the acronym with MIVF / SME before wiring"},
    {"name": "IUI", "color": AQUA,
     "purpose": "Intrauterine insemination — prepared sperm into uterus (no egg collection)",
     "steps": ["Consent & 3PC", "OI / Natural", "Folliculogram", "Trigger", "Semen Prep",
               "IUI Insemination", "Outcome"],
     "show": ["Overview", "Folliculogram", "Drugs", "IUI/DI", "Gametes", "Outcome Data", "Checklists"],
     "hide": ["Oocyte Retrieval", "Embryo Transfer", "Lab Results", "Stored Sample"],
     "fields": "Show partner & sperm source; hide # embryos to replace"},
    {"name": "Ovulation Induction", "color": MUTED,
     "purpose": "Drug-induced ovulation + timed intercourse (monitoring only)",
     "steps": ["Consent", "OI Drugs", "Folliculogram", "Trigger", "Timed Intercourse", "Outcome"],
     "show": ["Overview", "Drugs", "Folliculogram", "Outcome Data", "Checklists"],
     "hide": ["Oocyte Retrieval", "Embryo Transfer", "IUI/DI", "Gametes", "Lab Results", "Stored Sample"],
     "fields": "Mostly monitoring; hide embryo / sperm fields"},
    {"name": "Tracking", "color": INK,
     "purpose": "Cycle monitoring / scan tracking only (no procedure)",
     "steps": ["Folliculogram Tracking", "LH & Trigger Monitoring", "Outcome / Hand-off"],
     "show": ["Overview", "Folliculogram", "Outcome Data", "Drugs"],
     "hide": ["Oocyte Retrieval", "Embryo Transfer", "IUI/DI", "Gametes", "Lab Results",
              "Stored Sample", "Recommendation"],
     "fields": "Doctor, nurse, period dates only"},
    {"name": "Thaw and Refreeze", "color": RED,
     "purpose": "Thaw a stored specimen — (± biopsy) — refreeze — storage",
     "steps": ["Consent & 3PC", "Select Specimen", "Thaw", "Assess (± Biopsy)", "Refreeze",
               "Storage", "Witnessing"],
     "show": ["Overview", "Gametes", "Stored Sample", "Checklists"],
     "hide": ["Folliculogram", "Drugs", "Oocyte Retrieval", "Embryo Transfer", "IUI/DI",
              "Lab Results", "Outcome Data"],
     "fields": "Storage location & consent; hide clinical / stimulation fields"},
]

for idx, t in enumerate(TYPES):
    s = prs.slides.add_slide(BLANK); bg(s)
    header(s, t["name"].upper(), t["purpose"], 16 + idx)
    text(s, Inches(0.6), Inches(1.35), Inches(4), Inches(0.3), [[("KEY STEPS", 11, DEEP, True, HEAD)]])
    y2 = chip_flow(s, Inches(0.6), Inches(1.75), Inches(12.9), t["steps"], t["color"])
    top = y2 + Inches(0.4)
    ph = Inches(2.3)
    # SHOW panel
    box(s, Inches(0.6), top, Inches(6.0), ph, fill=WHITE, line=LINE, rounded=True)
    box(s, Inches(0.6), top, Inches(0.12), ph, fill=GREEN)
    text(s, Inches(0.85), top + Inches(0.15), Inches(5.5), Inches(0.4), [[("SHOW THESE TABS", 12, GREEN, True, HEAD)]])
    text(s, Inches(0.85), top + Inches(0.62), Inches(5.5), Inches(1.5), [[("  ·  ".join(t["show"]), 13, INK, False, BODY)]])
    # HIDE panel
    box(s, Inches(6.9), top, Inches(6.0), ph, fill=WHITE, line=LINE, rounded=True)
    box(s, Inches(6.9), top, Inches(0.12), ph, fill=RED)
    text(s, Inches(7.15), top + Inches(0.15), Inches(5.5), Inches(0.4), [[("HIDE THESE TABS", 12, RED, True, HEAD)]])
    text(s, Inches(7.15), top + Inches(0.62), Inches(5.5), Inches(1.5), [[("  ·  ".join(t["hide"]), 13, INK, False, BODY)]])
    # fields footer
    text(s, Inches(0.6), Inches(6.8), Inches(12.3), Inches(0.5),
         [[("Fields:  ", 12, DEEP, True, BODY), (t["fields"], 12, MUTED, False, BODY)]])

# ====================================================== SLIDE 24 — PLANNING TAB -> 3PC
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "TAB REDESIGN", "Treatment Planning  →  Three-Point Check", 24)
text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.8),
     [[("Plan first, then verify. ", 16, DEEP, True, BODY),
       ("The Treatment Planning tab consolidates today's Overview; the treatment type then hides "
        "irrelevant fields. The Three-Point Check gates every procedure.", 16, INK, False, BODY)]])
# planning card
box(s, Inches(0.6), Inches(2.5), Inches(7.4), Inches(4.1), fill=WHITE, line=LINE, rounded=True)
box(s, Inches(0.6), Inches(2.5), Inches(7.4), Inches(0.12), fill=AQUA)
text(s, Inches(0.85), Inches(2.7), Inches(7.0), Inches(0.4), [[("1 · TREATMENT PLANNING", 13, DEEP, True, HEAD)]])
plan = [
    ("Cycle Information", "type, template, billing, ANZARD, instructions"),
    ("People", "patient, partner, doctor, nurse, groups"),
    ("Clinic & Dates", "managing/procedure clinic, start & period dates"),
    ("Diagnosis", "female / male infertility + comment"),
    ("Plan (new)", "Event Template (seeds worklist) + Research Project"),
]
yy = Inches(3.2)
for t, d in plan:
    box(s, Inches(0.9), yy + Inches(0.06), Inches(0.14), Inches(0.14), fill=AQUA)
    text(s, Inches(1.2), yy - Inches(0.04), Inches(6.7), Inches(0.6),
         [[(t + "  ", 13, INK, True, BODY), (d, 12, MUTED, False, BODY)]])
    yy += Inches(0.62)
# 3pc card
box(s, Inches(8.2), Inches(2.5), Inches(4.7), Inches(4.1), fill=WHITE, line=LINE, rounded=True)
box(s, Inches(8.2), Inches(2.5), Inches(4.7), Inches(0.12), fill=RED)
text(s, Inches(8.45), Inches(2.7), Inches(4.2), Inches(0.4), [[("2 · THREE-POINT CHECK", 13, RED, True, HEAD)]])
tpc = ["Patient verified", "Partner verified", "Consent verified",
       "Independent witness ≠ performer", "Blocks procedure until Passed"]
yy = Inches(3.25)
for it in tpc:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), yy + Inches(0.02), Inches(0.26), Inches(0.26))
    c.fill.solid(); c.fill.fore_color.rgb = GREEN; c.line.fill.background(); c.shadow.inherit = False
    tf = c.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓"; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HEAD
    text(s, Inches(8.9), yy - Inches(0.03), Inches(3.9), Inches(0.5), [[(it, 12.5, INK, False, BODY)]])
    yy += Inches(0.6)
text(s, Inches(0.6), Inches(6.8), Inches(12.3), Inches(0.4),
     [[("Field-by-field content: OXar_Treatment_Cycle_Tab_Attributes.xlsx + 07-TREATMENT-PLANNING-TAB.md", 11.5, MUTED, False, BODY)]])

# ====================================================== SLIDE 25 — CRYO TANK MANAGEMENT
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "CRYO STORAGE", "Managing cryo tanks", 25)
text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.7),
     [[("Three linked screens — see capacity at a glance, drill into the position hierarchy, and "
        "add / edit tanks and slots (manager-gated).", 16, INK, False, BODY)]])
cards = [
    ("TANK WALL", "Card per tank with a fill gauge + status pill (OK / Low N2 / Alarm). + Add tank.", AQUA),
    ("TANK DETAIL", "Canister -> Goblet -> Straw/Vial tree; slots coloured by status & type; add/edit/retire.", DEEP),
    ("INVENTORY", "Who's in this tank: patient, specimen ID, material, freeze date, goblet/straw.", PURPLE),
]
cw = Inches(3.95); gap = Inches(0.2); x0 = Inches(0.6); y0 = Inches(2.4)
for i, (t, d, col) in enumerate(cards):
    x = x0 + i * (cw + gap)
    box(s, x, y0, cw, Inches(2.4), fill=WHITE, line=LINE, rounded=True)
    box(s, x, y0, cw, Inches(0.12), fill=col)
    chip(s, x + Inches(0.3), y0 + Inches(0.35), Inches(2.4), Inches(0.5), t, col, size=13)
    text(s, x + Inches(0.3), y0 + Inches(1.05), cw - Inches(0.6), Inches(1.2), [[(d, 13, INK, False, BODY)]])
box(s, Inches(0.6), Inches(5.15), Inches(12.3), Inches(1.55), fill=SOFT, rounded=True)
text(s, Inches(0.85), Inches(5.28), Inches(11.8), Inches(0.4),
     [[("Backend (all in the OXAR solution): ", 13, DEEP, True, BODY),
       ("CryoTank + CryoLocation full CRUD + CryoTank/GetInventory", 13, INK, True, MONO)]])
text(s, Inches(0.85), Inches(5.75), Inches(11.8), Inches(0.9),
     [[("Freeze picks a free slot → flips it Occupied; thaw frees it (→ Available), so the tank "
        "views always reflect live occupancy. Fill gauge: green <80%, amber 80-95%, red ≥95%.",
        12.5, MUTED, False, BODY)]])

prs.save("/workspace/oxarv1/docs/slides/OXar_Embryo_Module_Idiots_Guide.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")
